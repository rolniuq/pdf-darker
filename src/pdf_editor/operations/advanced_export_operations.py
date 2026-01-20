"""Advanced export operations for multiple formats."""

import os
import io
from pathlib import Path
from typing import Dict, List, Optional, Any
import csv
import json
import xml.etree.ElementTree as ET
from xml.dom import minidom
import fitz  # PyMuPDF
from PIL import Image
import pandas as pd
from docx import Document
from docx.shared import Inches
import xlsxwriter
import pptx
from pptx.util import Inches

from ..core.base import BaseOperation, ProcessingError, ValidationError
from ..utils.logging import get_logger

logger = get_logger("operations.advanced_export")


class ExportToWordOperation(BaseOperation):
    """Export PDF content to Word document."""
    
    def __init__(self, output_path: str, preserve_formatting: bool = True,
                 extract_images: bool = True, page_breaks: bool = True):
        super().__init__()
        self.output_path = Path(output_path)
        self.preserve_formatting = preserve_formatting
        self.extract_images = extract_images
        self.page_breaks = page_breaks
    
    def validate(self, document) -> None:
        """Validate export parameters."""
        if not self.output_path.suffix.lower() in ['.docx', '.doc']:
            raise ValidationError("Output file must have .docx or .doc extension")
    
    def execute(self, document) -> Dict:
        """Export PDF to Word document."""
        try:
            logger.info(f"Exporting PDF to Word: {self.output_path.name}")
            
            # Create Word document
            doc = Document()
            
            # Process each page
            total_pages = len(document)
            extracted_images = []
            
            for page_num in range(total_pages):
                page = document[page_num]
                
                # Extract and add text
                self._add_page_text_to_doc(doc, page, page_num + 1)
                
                # Extract and add images
                if self.extract_images:
                    page_images = self._extract_page_images(page, page_num)
                    extracted_images.extend(page_images)
                    
                    for img_info in page_images:
                        self._add_image_to_doc(doc, img_info)
                
                # Add page break
                if self.page_breaks and page_num < total_pages - 1:
                    doc.add_page_break()
            
            # Add document metadata
            self._add_document_metadata(doc, document)
            
            # Save document
            doc.save(str(self.output_path))
            
            file_size = self.output_path.stat().st_size
            
            logger.info(f"Successfully exported to Word: {file_size:,} bytes")
            
            return {
                'operation': 'export_to_word',
                'output_file': str(self.output_path),
                'pages_processed': total_pages,
                'images_extracted': len(extracted_images),
                'file_size': file_size,
                'preserve_formatting': self.preserve_formatting,
                'extract_images': self.extract_images,
                'page_breaks': self.page_breaks
            }
            
        except Exception as e:
            logger.error(f"Word export failed: {e}")
            raise ProcessingError(f"Word export failed: {e}")
    
    def _add_page_text_to_doc(self, doc: Document, page, page_num: int):
        """Add text from PDF page to Word document."""
        try:
            # Extract text with formatting information
            text_dict = page.get_text("dict")
            
            # Add page header
            if page_num > 1:
                doc.add_heading(f'Page {page_num}', level=2)
            
            # Process text blocks
            for block in text_dict.get("blocks", []):
                if "lines" in block:
                    for line in block["lines"]:
                        for span in line["spans"]:
                            # Add text with basic formatting
                            paragraph = doc.add_paragraph(span["text"])
                            
                            # Apply basic formatting (simplified)
                            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run("")
                            
                            if span.get("bold"):
                                run.bold = True
                            if span.get("italic"):
                                run.italic = True
                            
                            # Set font size
                            if "size" in span:
                                run.font.size = span["size"] * 0.75  # Convert to points
                
                elif "image" in block:
                    # Handle image blocks (already handled separately)
                    pass
        
        except Exception as e:
            logger.warning(f"Failed to extract text from page {page_num}: {e}")
    
    def _extract_page_images(self, page, page_num: int) -> List[Dict]:
        """Extract images from PDF page."""
        images = []
        image_list = page.get_images()
        
        for img_index, img in enumerate(image_list):
            try:
                xref = img[0]
                pix = fitz.Pixmap(page.parent, xref)
                
                if pix.n - pix.alpha < 4:  # Ensure it's not a mask or CMYK
                    img_data = pix.tobytes("png")
                    
                    images.append({
                        'page': page_num,
                        'index': img_index,
                        'data': img_data,
                        'width': pix.width,
                        'height': pix.height
                    })
                
                pix = None  # Free memory
                
            except Exception as e:
                logger.warning(f"Failed to extract image {img_index} from page {page_num}: {e}")
        
        return images
    
    def _add_image_to_doc(self, doc: Document, img_info: Dict):
        """Add extracted image to Word document."""
        try:
            # Save image data to temporary file
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                tmp_file.write(img_info['data'])
                tmp_file_path = tmp_file.name
            
            # Add image to document with proper sizing
            max_width = Inches(6)  # Maximum width of 6 inches
            max_height = Inches(4)  # Maximum height of 4 inches
            
            # Calculate dimensions to maintain aspect ratio
            width_inches = img_info['width'] / 96  # Assuming 96 DPI
            height_inches = img_info['height'] / 96
            
            # Scale to fit within maximum dimensions
            if width_inches > max_width:
                height_inches = (max_width / width_inches) * height_inches
                width_inches = max_width
            
            if height_inches > max_height:
                width_inches = (max_height / height_inches) * width_inches
                height_inches = max_height
            
            # Add image to document
            doc.add_picture(tmp_file_path, width=width_inches, height=height_inches)
            
            # Clean up temporary file
            os.unlink(tmp_file_path)
            
        except Exception as e:
            logger.warning(f"Failed to add image to Word document: {e}")
    
    def _add_document_metadata(self, doc: Document, document):
        """Add metadata to Word document."""
        try:
            metadata = document.metadata
            
            if metadata.get('title'):
                doc.core_properties.title = metadata['title']
            
            if metadata.get('author'):
                doc.core_properties.author = metadata['author']
            
            if metadata.get('subject'):
                doc.core_properties.subject = metadata['subject']
            
            if metadata.get('keywords'):
                doc.core_properties.keywords = metadata['keywords']
            
        except Exception as e:
            logger.warning(f"Failed to add document metadata: {e}")


class ExportToExcelOperation(BaseOperation):
    """Export PDF form data or structured content to Excel."""
    
    def __init__(self, output_path: str, export_type: str = 'form_data',
                 include_metadata: bool = True):
        super().__init__()
        self.output_path = Path(output_path)
        self.export_type = export_type.lower()  # 'form_data', 'table_data', 'text_blocks'
        self.include_metadata = include_metadata
    
    def validate(self, document) -> None:
        """Validate export parameters."""
        if not self.output_path.suffix.lower() in ['.xlsx', '.xls']:
            raise ValidationError("Output file must have .xlsx or .xls extension")
        
        if self.export_type not in ['form_data', 'table_data', 'text_blocks']:
            raise ValidationError("Export type must be 'form_data', 'table_data', or 'text_blocks'")
    
    def execute(self, document) -> Dict:
        """Export PDF to Excel workbook."""
        try:
            logger.info(f"Exporting PDF to Excel: {self.output_path.name}")
            
            # Create Excel workbook
            workbook = xlsxwriter.Workbook(str(self.output_path))
            
            # Add formats
            header_format = workbook.add_format({
                'bold': True,
                'bg_color': '#D7E4BC',
                'border': 1
            })
            
            data_format = workbook.add_format({
                'border': 1,
                'text_wrap': True
            })
            
            # Export based on type
            if self.export_type == 'form_data':
                self._export_form_data(workbook, document, header_format, data_format)
            elif self.export_type == 'table_data':
                self._export_table_data(workbook, document, header_format, data_format)
            elif self.export_type == 'text_blocks':
                self._export_text_blocks(workbook, document, header_format, data_format)
            
            # Close workbook
            workbook.close()
            
            file_size = self.output_path.stat().st_size
            
            logger.info(f"Successfully exported to Excel: {file_size:,} bytes")
            
            return {
                'operation': 'export_to_excel',
                'output_file': str(self.output_path),
                'export_type': self.export_type,
                'file_size': file_size,
                'include_metadata': self.include_metadata
            }
            
        except Exception as e:
            logger.error(f"Excel export failed: {e}")
            raise ProcessingError(f"Excel export failed: {e}")
    
    def _export_form_data(self, workbook, document, header_format, data_format):
        """Export form field data to Excel."""
        worksheet = workbook.add_worksheet('Form Data')
        
        # Headers
        headers = ['Field Name', 'Field Type', 'Value', 'Page', 'Position (X,Y)', 'Flags']
        for col, header in enumerate(headers):
            worksheet.write(0, col, header, header_format)
        
        # Extract form data
        row = 1
        for page_num in range(len(document)):
            page = document[page_num]
            widgets = page.widgets()
            
            for widget in widgets:
                field_name = widget.field_name or f"Field_{row}"
                field_type = widget.field_type or "unknown"
                field_value = widget.field_value or ""
                
                # Position
                rect = widget.rect
                position = f"({rect.x0:.0f}, {rect.y0:.0f})"
                
                # Flags
                flags = []
                if widget.is_readonly:
                    flags.append("Readonly")
                if widget.is_required:
                    flags.append("Required")
                
                worksheet.write(row, 0, field_name, data_format)
                worksheet.write(row, 1, field_type, data_format)
                worksheet.write(row, 2, field_value, data_format)
                worksheet.write(row, 3, page_num + 1, data_format)
                worksheet.write(row, 4, position, data_format)
                worksheet.write(row, 5, ", ".join(flags), data_format)
                
                row += 1
        
        # Auto-adjust column widths
        for col in range(len(headers)):
            worksheet.set_column(col, col, max(len(header) for header in headers) + 5)
    
    def _export_table_data(self, workbook, document, header_format, data_format):
        """Extract and export table data to Excel."""
        worksheet = workbook.add_worksheet('Table Data')
        
        # This is a simplified implementation
        # In practice, you'd need sophisticated table detection
        row = 0
        
        for page_num in range(len(document)):
            page = document[page_num]
            
            # Look for table-like structures (simplified)
            text_lines = page.get_text().split('\n')
            
            current_table = []
            for line in text_lines:
                if line.strip():
                    # Simple table detection (tabs or multiple spaces)
                    if '\t' in line or '  ' in line:
                        cells = [cell.strip() for cell in line.replace('\t', '  ').split('  ') if cell.strip()]
                        current_table.append(cells)
                    elif current_table:
                        # End of table
                        if len(current_table) > 1:
                            # Write table to Excel
                            for table_row in current_table:
                                for col, cell in enumerate(table_row):
                                    worksheet.write(row, col, cell, data_format)
                                row += 1
                            row += 1  # Add space between tables
                        current_table = []
        
        # Write any remaining table
        if current_table:
            for table_row in current_table:
                for col, cell in enumerate(table_row):
                    worksheet.write(row, col, cell, data_format)
                row += 1
    
    def _export_text_blocks(self, workbook, document, header_format, data_format):
        """Export text blocks to Excel."""
        worksheet = workbook.add_worksheet('Text Blocks')
        
        # Headers
        headers = ['Page', 'Block', 'Text', 'Font', 'Size', 'Position (X,Y)', 'Type']
        for col, header in enumerate(headers):
            worksheet.write(0, col, header, header_format)
        
        row = 1
        block_num = 1
        
        for page_num in range(len(document)):
            page = document[page_num]
            text_dict = page.get_text("dict")
            
            for block in text_dict.get("blocks", []):
                if "lines" in block:
                    for line in block["lines"]:
                        for span in line["spans"]:
                            text = span["text"]
                            font = span.get("font", "unknown")
                            size = span.get("size", 12)
                            x, y = span.get("origin", (0, 0))
                            block_type = "text"
                            
                            worksheet.write(row, 0, page_num + 1, data_format)
                            worksheet.write(row, 1, block_num, data_format)
                            worksheet.write(row, 2, text, data_format)
                            worksheet.write(row, 3, font, data_format)
                            worksheet.write(row, 4, size, data_format)
                            worksheet.write(row, 5, f"({x:.0f}, {y:.0f})", data_format)
                            worksheet.write(row, 6, block_type, data_format)
                            
                            row += 1
                            block_num += 1


class ExportToPowerPointOperation(BaseOperation):
    """Export PDF content to PowerPoint presentation."""
    
    def __init__(self, output_path: str, one_slide_per_page: bool = True,
                 slide_size: str = 'standard_4_3', extract_images: bool = True):
        super().__init__()
        self.output_path = Path(output_path)
        self.one_slide_per_page = one_slide_per_page
        self.slide_size = slide_size
        self.extract_images = extract_images
    
    def validate(self, document) -> None:
        """Validate export parameters."""
        if not self.output_path.suffix.lower() in ['.pptx', '.ppt']:
            raise ValidationError("Output file must have .pptx or .ppt extension")
        
        valid_sizes = ['standard_4_3', 'standard_16_9', 'widescreen', 'a4']
        if self.slide_size not in valid_sizes:
            raise ValidationError(f"Slide size must be one of: {', '.join(valid_sizes)}")
    
    def execute(self, document) -> Dict:
        """Export PDF to PowerPoint presentation."""
        try:
            logger.info(f"Exporting PDF to PowerPoint: {self.output_path.name}")
            
            # Create PowerPoint presentation
            prs = pptx.Presentation()
            
            # Set slide size
            self._set_slide_size(prs)
            
            # Process pages
            total_pages = len(document)
            slides_created = 0
            images_extracted = 0
            
            if self.one_slide_per_page:
                # One slide per page
                for page_num in range(total_pages):
                    page = document[page_num]
                    
                    # Create slide from page
                    slide, slide_images = self._create_slide_from_page(prs, page, page_num + 1)
                    
                    if slide:
                        slides_created += 1
                        images_extracted += len(slide_images)
            else:
                # Combine multiple pages per slide
                self._create_slides_from_pages(prs, document)
                slides_created = len(prs.slides)
            
            # Save presentation
            prs.save(str(self.output_path))
            
            file_size = self.output_path.stat().st_size
            
            logger.info(f"Successfully exported to PowerPoint: {file_size:,} bytes")
            
            return {
                'operation': 'export_to_powerpoint',
                'output_file': str(self.output_path),
                'pages_processed': total_pages,
                'slides_created': slides_created,
                'images_extracted': images_extracted,
                'file_size': file_size,
                'one_slide_per_page': self.one_slide_per_page,
                'slide_size': self.slide_size
            }
            
        except Exception as e:
            logger.error(f"PowerPoint export failed: {e}")
            raise ProcessingError(f"PowerPoint export failed: {e}")
    
    def _set_slide_size(self, prs):
        """Set PowerPoint slide size."""
        from pptx.util import Inches
        
        size_map = {
            'standard_4_3': (Inches(10), Inches(7.5)),
            'standard_16_9': (Inches(10), Inches(5.625)),
            'widescreen': (Inches(13.33), Inches(7.5)),
            'a4': (Inches(10), Inches(7.07))
        }
        
        width, height = size_map.get(self.slide_size, size_map['standard_4_3'])
        prs.slide_width = width
        prs.slide_height = height
    
    def _create_slide_from_page(self, prs, page, page_num: int):
        """Create a PowerPoint slide from a PDF page."""
        try:
            # Convert page to image
            pix = page.get_pixmap(dpi=150)
            img_data = pix.tobytes("png")
            
            # Save to temporary file
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                tmp_file.write(img_data)
                tmp_file_path = tmp_file.name
            
            # Create slide
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add image to slide (fit to slide)
            slide.shapes.add_picture(
                tmp_file_path,
                0, 0,
                width=prs.slide_width,
                height=prs.slide_height
            )
            
            # Add page number as text
            from pptx.util import Pt
            left = Inches(0.5)
            top = prs.slide_height - Inches(0.5)
            width = Inches(2)
            height = Inches(0.5)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = f"Page {page_num}"
            p.font.size = Pt(12)
            
            # Clean up
            os.unlink(tmp_file_path)
            
            return slide, []
            
        except Exception as e:
            logger.warning(f"Failed to create slide from page {page_num}: {e}")
            return None, []
    
    def _create_slides_from_pages(self, prs, document):
        """Create slides by combining multiple PDF pages."""
        # This is a simplified implementation
        # In practice, you'd want sophisticated content analysis
        pass


# Required imports
import tempfile